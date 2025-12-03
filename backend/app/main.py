import sys
import os
import json
from pprint import pprint
from pptx import Presentation
 
# Ensure project root is in path
sys.path.append(os.getcwd())
 
try:
    # Import 4 agents
    from backend.app.agents.requirement_agent import requirement_node
    from backend.app.agents.retrieval_agent import retrieval_node
    from backend.app.agents.comparator_agent import comparator_node
    from backend.app.agents.pitch_agent import sales_pitch_node
   
    # Import PPT generator agent
    from backend.app.agents.deck_agent import create_professional_ppt
 
    # State + LLM
    from backend.app.graph.state import AgentState
    from backend.app.llm.llm_client import llm
 
except ImportError as e:
    print(f"❌ Import Error: {e}")
    sys.exit(1)
 
 
# ==========================
# Utility: Save JSON
# ==========================
def save_json(path, data):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)
    print(f"💾 Saved: {path}")
 
 
def extract_ppt_text(ppt_path):
    if not ppt_path or not os.path.exists(ppt_path):
        return ""
 
    prs = Presentation(ppt_path)
    collected = []
 
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text:
                    collected.append(shape.text.strip())
            except:
                pass
 
    return "\n\n".join(collected)
 
 
# For nice log sections
def print_section(title):
    print(f"\n{'='*65}")
    print(f" 🕵️  {title}")
    print(f"{'='*65}")
 
 
# ============================================================
#                 FULL PIPELINE (WITH JSON SAVING)
# ============================================================
def test_full_pipeline():
    query = "We need 10 laptops for a design team, around ₹1 lakh each, lightweight, 8GB RAM, and good graphics."
 
    # DIRECTORIES
    EVAL_DIR = "evaluation_outputs"
    os.makedirs(EVAL_DIR, exist_ok=True)
 
    # INITIAL STATE
    state: AgentState = {
        "user_query": query,
        "requirements": {},
        "retrieved_products": [],
        "ranked_products": [],
        "sales_pitch": {},
        "deck_file_path": None,
        "llm": llm
    }
 
    # ---------------------------------------------------------------
    # 1️⃣  REQUIREMENT AGENT
    # ---------------------------------------------------------------
    print_section("STEP 1: REQUIREMENT AGENT")
    update1 = requirement_node(state)
    state.update(update1)
 
    save_json(f"{EVAL_DIR}/requirements.json", state["requirements"])
 
    # ---------------------------------------------------------------
    # 2️⃣  RETRIEVAL AGENT
    # ---------------------------------------------------------------
    print_section("STEP 2: RETRIEVAL AGENT")
    update2 = retrieval_node(state)
    state.update(update2)
 
    save_json(f"{EVAL_DIR}/retrieved_products.json", state["retrieved_products"])
 
    # ---------------------------------------------------------------
    # 3️⃣  COMPARATOR AGENT
    # ---------------------------------------------------------------
    print_section("STEP 3: COMPARATOR AGENT")
    update3 = comparator_node(state)
    state.update(update3)
 
    save_json(f"{EVAL_DIR}/ranked_products.json", state["ranked_products"])
 
    # ---------------------------------------------------------------
    # 4️⃣  SALES PITCH AGENT
    # ---------------------------------------------------------------
    print_section("STEP 4: SALES PITCH AGENT")
    update4 = sales_pitch_node(state)
    state.update(update4)
 
    save_json(f"{EVAL_DIR}/sales_pitch.json", state["sales_pitch"])
 
    # ---------------------------------------------------------------
    # 5️⃣  PPT GENERATOR AGENT
    # ---------------------------------------------------------------
    print_section("STEP 5: PPT GENERATION")
 
    deck_path = create_professional_ppt(
        sales_pitch=state["sales_pitch"],
        products=state["ranked_products"][:5],
        customer_query=query,
        output_file="HP_Proposal_Professional.pptx"
    )
 
    state["deck_file_path"] = deck_path
 
    print(f"📁 PPT generated: {deck_path}")
 
    # Extract PPT text for RAGAS
    ppt_text = extract_ppt_text(deck_path)
    save_json(f"{EVAL_DIR}/ppt_text.json", {"ppt_text": ppt_text})
 
    # ---------------------------------------------------------------
    # 6️⃣  FINAL SUMMARY FILE
    # ---------------------------------------------------------------
    summary = {
        "user_query": query,
        "ppt_path": deck_path,
        "requirements_file": f"{EVAL_DIR}/requirements.json",
        "retrieved_products_file": f"{EVAL_DIR}/retrieved_products.json",
        "ranked_products_file": f"{EVAL_DIR}/ranked_products.json",
        "sales_pitch_file": f"{EVAL_DIR}/sales_pitch.json",
        "ppt_text_file": f"{EVAL_DIR}/ppt_text.json",
    }
 
    save_json(f"{EVAL_DIR}/pipeline_summary.json", summary)
 
    print_section("PIPELINE COMPLETED SUCCESSFULLY")
    print("🎉 All agent outputs stored for RAGAS evaluation.")
 
 
if __name__ == "__main__":
    test_full_pipeline()