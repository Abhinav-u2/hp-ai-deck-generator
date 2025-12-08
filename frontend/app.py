import streamlit as st
import requests
import io
import base64
import os
import json  # 👈 NEW: for parsing streaming JSON lines
from PIL import Image
from pptx import Presentation
from dotenv import load_dotenv

# Load Environment Variables
load_dotenv()

# Configuration
FASTAPI_BASE_URL = os.getenv("FASTAPI_BASE_URL", "http://localhost:8000").rstrip("/")

# ==========================================
# 🎨 GLOBAL THEME & CUSTOM CSS
# ==========================================
HP_BLUE = "#004A8F"
HP_BLUE_LIGHT = "#0096D6"
HP_BLUE_SOFT = "#E6F0FA"

st.set_page_config(
    page_title="HP AI Sales Deck Generator",
    layout="wide",
    page_icon="💻"
)

st.markdown(
    f"""
    <style>
    /* Global page */
    .stApp {{
        background: linear-gradient(180deg, #f5f8fc 0%, #eef3fa 45%, #f7faff 100%);
        color: #1f2a3c;
        font-family: "Segoe UI", system-ui, -apple-system, BlinkMacSystemFont, sans-serif;
    }}

    /* Remove extra padding at top */
    .block-container {{
        padding-top: 1.2rem !important;
        padding-bottom: 1.5rem !important;
    }}

    /* Main header */
    .app-header {{
        padding: 1.1rem 1.2rem;
        border-radius: 18px;
        background: linear-gradient(120deg, {HP_BLUE}, {HP_BLUE_LIGHT});
        box-shadow: 0 14px 30px rgba(0,0,0,0.25);
        color: white;
        position: relative;
        overflow: hidden;
    }}

    .app-header::before {{
        content: "";
        position: absolute;
        top: -40px;
        right: -80px;
        width: 220px;
        height: 220px;
        background: radial-gradient(circle, rgba(255,255,255,0.35), transparent 60%);
        opacity: 0.9;
        filter: blur(2px);
    }}

    .app-header h1 {{
        font-size: 1.9rem;
        margin-bottom: 0.2rem;
        position: relative;
        z-index: 2;
    }}
    .app-header p {{
        margin-top: 0;
        opacity: 0.95;
        position: relative;
        z-index: 2;
    }}

    /* Step cards */
    .step-card {{
        border-radius: 18px;
        padding: 1rem 1.2rem;
        margin-bottom: 1rem;
        background: #ffffff;
        border: 1px solid rgba(0, 87, 174, 0.18);
        box-shadow: 0 10px 24px rgba(5,36,80,0.12);
        position: relative;
        overflow: hidden;
        animation: slideInUp 0.5s ease-out;
    }}

    .step-header {{
        display: flex;
        align-items: center;
        gap: 0.7rem;
        margin-bottom: 0.45rem;
    }}

    .step-badge {{
        width: 32px;
        height: 32px;
        border-radius: 999px;
        background: radial-gradient(circle at 30% 20%, #ffffff, #cde5ff);
        display: flex;
        align-items: center;
        justify-content: center;
        font-weight: 700;
        color: {HP_BLUE};
        box-shadow: 0 0 14px rgba(0,84,175,0.45);
        flex-shrink: 0;
    }}

    .step-title {{
        font-weight: 600;
        font-size: 1.02rem;
        color: #12233a;
    }}

    .step-subtitle {{
        font-size: 0.85rem;
        opacity: 0.8;
    }}

    .step-status {{
        font-size: 0.8rem;
        padding: 0.2rem 0.65rem;
        border-radius: 999px;
        border: 1px solid rgba(0, 153, 255, 0.7);
        color: #006bb3;
        display: inline-flex;
        align-items: center;
        gap: 0.2rem;
        margin-left: auto;
        background: rgba(230, 243, 255, 0.9);
    }}

    .pill-title {{
        font-size: 0.8rem;
        text-transform: uppercase;
        letter-spacing: 0.08em;
        opacity: 0.65;
        margin-bottom: 0.35rem;
        color: #4a5c73;
    }}

    .highlight-box {{
        border-radius: 12px;
        padding: 0.65rem 0.8rem;
        margin-bottom: 0.6rem;
        background: #f7fbff;
        border: 1px solid rgba(0, 125, 255, 0.22);
    }}
    .highlight-box h4 {{
        margin: 0 0 0.3rem 0;
        font-size: 0.96rem;
        color: #11345b;
    }}
    .highlight-box ul {{
        padding-left: 1rem;
        margin: 0.2rem 0;
    }}
    .highlight-box li {{
        font-size: 0.86rem;
        margin-bottom: 0.18rem;
    }}

    /* PPT preview styling */
    .slide-preview {{
        border-radius: 18px;
        padding: 0.8rem 0.9rem 0.9rem 0.9rem;
        background: #ffffff;
        border: 1px solid rgba(0, 87, 174, 0.18);
        box-shadow: 0 10px 24px rgba(5,36,80,0.12);
    }}

    .preview-title {{
        font-size: 0.9rem;
        opacity: 0.85;
        margin-bottom: 0.4rem;
        color: #183d66;
    }}

    /* Download button */
    .download-btn a {{
        text-decoration: none;
        color: white !important;
    }}

    .hp-btn {{
        background: linear-gradient(120deg, {HP_BLUE_LIGHT}, {HP_BLUE});
        border-radius: 999px;
        border: none;
        padding: 0.65rem 1.4rem;
        color: white;
        font-weight: 600;
        font-size: 0.9rem;
        cursor: pointer;
        box-shadow: 0 14px 24px rgba(0,0,0,0.25);
        margin-top: 0.4rem;
    }}

    /* Slide-in animation */
    @keyframes slideInUp {{
        0% {{
            opacity: 0;
            transform: translateY(18px);
        }}
        100% {{
            opacity: 1;
            transform: translateY(0);
        }}
    }}
    </style>
    """,
    unsafe_allow_html=True
)

# ==========================================
# 🛠️ HELPER: PPTX TO IMAGES
# ==========================================
def pptx_to_images(ppt_bytes, dpi=96):
    """
    Render slides as simple preview images (title + bullet text).
    """
    try:
        prs = Presentation(io.BytesIO(ppt_bytes))
        images = []
        width_px = int(prs.slide_width / 914400 * dpi)
        height_px = int(prs.slide_height / 914400 * dpi)

        for slide in prs.slides:
            img = Image.new("RGB", (width_px, height_px), "white")
            from PIL import ImageDraw, ImageFont
            draw = ImageDraw.Draw(img)

            # Choose fonts
            try:
                title_font = ImageFont.truetype("arial.ttf", 32)
                body_font = ImageFont.truetype("arial.ttf", 18)
            except Exception:
                title_font = None
                body_font = None

            y = 40
            # Title
            if slide.shapes.title and slide.shapes.title.text:
                title_text = slide.shapes.title.text.strip()
                draw.text((40, y), title_text[:60], fill=(0, 55, 110), font=title_font)
                y += 60

            # Limited text content from slide
            displayed_lines = 0
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    if shape == slide.shapes.title:
                        continue
                    text = shape.text.strip()
                    if text:
                        for line in text.splitlines():
                            if displayed_lines > 10:
                                break
                            line = line.strip()
                            if not line:
                                continue
                            draw.text((60, y), f"• {line[:90]}", fill=(30, 30, 30), font=body_font)
                            y += 28
                            displayed_lines += 1

            images.append(img)

        return images
    except Exception as e:
        st.error(f"Preview Render Error: {e}")
        return []


def get_download_link(bytes_obj, filename="HP_Proposal_Professional.pptx"):
    b64 = base64.b64encode(bytes_obj).decode()
    return f"data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64}"



def parse_sales_pitch_blocks(raw_text: str):
    """
    Safer parser for LLM output.
    Returns [{name: "", bullets: []}]
    """
    if not raw_text:
        return []

    blocks = [b.strip() for b in raw_text.split("\n") if b.strip()]
    parsed = []

    for block in blocks:
        if " - " in block:
            first, *rest = block.split(" - ")
            product = first.replace("**", "").strip()

            bullets = [
                r.strip()
                .replace("</ul>", "")
                .replace("</div>", "")
                .replace("<ul>", "")
                .replace("<div>", "")
                for r in rest
                if r.strip()
            ]

            parsed.append({"name": product, "bullets": bullets})

    return parsed



def render_product_card(title, bullet_list, color="#004A8F"):
    bullet_html = "".join([
        f"<li style='margin-bottom:4px;'>{b.replace('</ul>', '').replace('</div>', '').replace('<ul>', '').replace('<div>', '')}</li>"
        for b in bullet_list
    ])


    st.markdown(
        f"""
        <div style="
            border: 1px solid rgba(0,87,174,0.25);
            border-left: 6px solid {color};
            padding: 14px 16px;
            margin: 8px 0;
            border-radius: 12px;
            background: linear-gradient(145deg, #ffffff, #f6faff);
            box-shadow: 0 6px 14px rgba(0,0,0,0.08);
        ">
            <h4 style="margin:0; color:{color};">{title}</h4>
            <ul style="padding-left:18px; margin-top:8px; font-size:14.5px; line-height:1.45;">
                {bullet_html}
            </ul>
        </div>
        """,
        unsafe_allow_html=True
    )

# ==========================================
# 🧠 HEADER
# ==========================================
st.markdown(
    """
    <div class="app-header">
        <h1>💻 HP AI Sales Deck Generator</h1>
        <p>Turn a natural-language customer requirement into an HP-branded sales proposal deck.</p>
    </div>
    """,
    unsafe_allow_html=True
)

st.markdown("")


# ==========================================
# ✍️ INPUT SECTION
# ==========================================
with st.container():
    col_input, col_side = st.columns([2.4, 1.6], gap="large")

    with col_input:
        st.markdown("#### 1️⃣ Enter Customer Requirement")
        query = st.text_area(
            "Customer Requirement",
            height=100,
            placeholder="E.g., We need 10 laptops for a design team, around ₹1 lakh each, lightweight, 8GB RAM, and good graphics.",
            label_visibility="collapsed"
        )
        run_btn = st.button("⚡ Generate HP Sales Deck", type="primary")

    with col_side:
        st.markdown("#### 🎯 What this app does")
        st.markdown(
            """
            - **Requirement Agent** understands the customer need  
            - **Hybrid Retrieval Agent** searches HP catalogue  
            - **Comparator Agent** ranks matching HP products  
            - **Sales Pitch Agent** creates advantages & reasons to buy  
            - **Deck Agent** exports an HP-branded PowerPoint  
            """
        )

def score_bar(score):
    percentage = (score / 10) * 100
    if score >= 8:
        color = "#4CAF50"   
    elif score >= 5:
        color = "#FFC107"   
    else:
        color = "#F44336"

    return f"""
    <div style="margin-top:3px; width:100%; background:#eee; border-radius:8px;">
        <div style="
            width:{percentage}%;
            background:{color};
            padding:5px;
            border-radius:8px;
            text-align:center;
            color:white;
            font-weight:bold;
            font-size:0.85rem;
        ">{score}/10</div>
    </div>
    """


# ==========================================
# 🚀 MAIN EXECUTION (STREAMING VIA /query-stream)
# ==========================================
if run_btn and not query:
    st.warning("Please enter a requirement to start.")

if run_btn and query:
    st.divider()

    # Initialize containers for agent outputs so they exist even if something fails
    requirements = {}
    retrieved_products = []
    ranked_products = []
    sales_pitch = {}
    product_highlights = ""
    reasons_to_buy = ""
    ppt_bytes = b""

    with st.status("🤖 Running HP AI agent workflow...", expanded=True) as status:
        try:
            # STEP 1: /query-stream (stream the pipeline step by step)
            status.write("🧠 Step 1 – Connecting to /query-stream and running agents sequentially...")

            stream_url = f"{FASTAPI_BASE_URL}/query-stream"
            payload = {"query": query}

            with requests.post(stream_url, json=payload, stream=True, timeout=300) as resp:
                resp.raise_for_status()
                status.write("🔌 Connected. Waiting for agent updates...")

                for raw_line in resp.iter_lines():
                    if not raw_line:
                        continue

                    decoded = raw_line.decode("utf-8").strip()
                    if not decoded:
                        continue

                    if decoded == "END":
                        # Backend signaled completion
                        status.write("✅ All pipeline steps completed on backend.")
                        break

                    # Each non-END line is a JSON event
                    try:
                        event = json.loads(decoded)
                    except json.JSONDecodeError:
                        # Ignore malformed lines
                        continue

                    step = event.get("step")
                    label = event.get("label")
                    data = event.get("data")

                    if step == 1 and label == "requirements":
                        requirements = data or {}
                        status.write("✅ Step 1 – Requirement identified according to query.")

                    elif step == 2 and label == "retrieved_products":
                        retrieved_products = data or []
                        status.write(
                            f"✅ Step 2 – Fetched the products."
                        )

                    elif step == 3 and label == "ranked_products":
                        ranked_products = data or []
                        status.write(
                            f"✅ Step 3 – Ranked the products."
                        )

                    elif step == 4 and label == "sales_pitch":
                        sales_pitch = data or {}
                        status.write("✅ Step 4 – Sales Pitch generated.")

                    elif step == 5 and label == "ppt_ready":
                        status.write("📁 Step 5 – Deck finished generating the PPT.")

            # Extract pitch strings after streaming is done
            product_highlights = sales_pitch.get("product_highlights", "") if isinstance(sales_pitch, dict) else ""
            reasons_to_buy = sales_pitch.get("reasons_to_buy", "") if isinstance(sales_pitch, dict) else ""

            # STEP 2: /get_ppt (actual PPT file once backend signals ppt_ready)
            status.write("📊 Fetching generated PowerPoint from /get_ppt...")

            ppt_url = f"{FASTAPI_BASE_URL}/get_ppt"
            file_resp = requests.get(ppt_url, params={"query": query}, timeout=300)
            file_resp.raise_for_status()
            ppt_bytes = file_resp.content

            # STEP 3: /get_pdf (Preview / Download PDF)
            status.write("📄 Attempting to retrieve PDF preview...")

            pdf_bytes = None
            try:
                pdf_url = f"{FASTAPI_BASE_URL}/get_pdf"
                pdf_resp = requests.get(pdf_url, params={"query": query}, timeout=200)
                pdf_resp.raise_for_status()
                pdf_bytes = pdf_resp.content
                status.write("📃 PDF preview retrieved successfully.")
            except Exception:
                status.write("⚠ PDF preview not available (conversion may be disabled).")

            status.write("📁 PPT retrieved successfully.")
            status.update(label="✅ HP AI workflow complete!", state="complete", expanded=False)

        except requests.exceptions.RequestException as e:
            st.error(f"Server Connection Error: {e}")
            status.update(label="❌ Failed", state="error")
            st.stop()
        except Exception as e:
            st.error(f"Unexpected Error: {e}")
            status.update(label="❌ Error", state="error")
            st.stop()

    # tiny celebration
    st.balloons()

    # ==========================================
    # 📚 VERTICAL STORY: WORKFLOW EXPLAINABILITY
    # ==========================================
    st.markdown("## 🔄 AI Workflow – From Requirement to HP Sales Deck")

    # ---- Step 1: Requirement Agent ----
    with st.container():
        st.markdown(
            """
            <div class="step-card">
                <div class="step-header">
                    <div class="step-badge">1</div>
                    <div>
                        <div class="step-title">Requirement Agent – Understands the Customer Need</div>
                        <div class="step-subtitle">Parses free-form input into a structured requirement for HP products.</div>
                    </div>
                    <div class="step-status">✅ Completed</div>
                </div>
            """,
            unsafe_allow_html=True
        )

        col1, col2 = st.columns([1.3, 1.7])
        with col1:
            st.markdown('<div class="pill-title">Extracted Requirement</div>', unsafe_allow_html=True)
            if isinstance(requirements, dict) and requirements:
                st.markdown(
                    f"""
                    <div class="highlight-box">
                        <h4>Customer Need</h4>
                        <ul>
                            <li><b>Category:</b> {requirements.get('product_category', '—')}</li>
                            <li><b>Quantity:</b> {requirements.get('quantity', '—')}</li>
                            <li><b>Budget / Unit:</b> {requirements.get('budget_per_unit', '—')}</li>
                            <li><b>User Intent:</b> {requirements.get('user_intent', '—')}</li>
                        </ul>
                    </div>
                    """,
                    unsafe_allow_html=True
                )

                spec_req = requirements.get("specific_requirements", {})
                if spec_req:
                    bullets = []
                    for k, v in spec_req.items():
                        bullets.append(f"- **{k.capitalize()}**: {v}")
                    st.markdown(
                        "<div class='pill-title'>Specific Requirements</div>",
                        unsafe_allow_html=True
                    )
                    st.markdown(
                        "<div class='highlight-box'>" + "<br>".join(bullets) + "</div>",
                        unsafe_allow_html=True
                    )
            else:
                st.info("No structured requirements returned from backend.")

        with col2:
            st.markdown('<div class="pill-title">How this helps</div>', unsafe_allow_html=True)
            st.markdown(
                """
                - Converts vague text into **clear product specs**  
                - Guides downstream agents towards the right **HP series**  
                - Captures constraints like **RAM, weight, graphics, budget**  
                """
            )
        st.markdown("</div>", unsafe_allow_html=True)

    # ---- Step 2: Retrieval Agent ----
    with st.container():
        st.markdown(
            """
            <div class="step-card">
                <div class="step-header">
                    <div class="step-badge">2</div>
                    <div>
                        <div class="step-title">Hybrid Retrieval Agent – Searches HP Catalogue</div>
                        <div class="step-subtitle">Runs dense + sparse hybrid search over HP product vectors.</div>
                    </div>
                    <div class="step-status">✅ Completed</div>
                </div>
            """,
            unsafe_allow_html=True
        )

    # ---- Step 3: Comparator Agent ----
    with st.container():
        st.markdown(
            """
            <div class="step-card">
                <div class="step-header">
                    <div class="step-badge">3</div>
                    <div>
                        <div class="step-title">Comparator Agent – Ranks HP Products</div>
                        <div class="step-subtitle">Evaluates RAM, graphics, weight & fit to the customer scenario.</div>
                    </div>
                    <div class="step-status">✅ Completed</div>
                </div>
            """,
            unsafe_allow_html=True
        )

        if ranked_products:
            for idx, p in enumerate(ranked_products[:5], start=1):
                st.markdown(
                    f"""
                    <div class="highlight-box">
                        <h4>#{idx} – {p.get('product_name', 'Unknown')}</h4>
                        <ul>
                            <li><b>Match Score:</b> {round(p.get('score', 0.0), 2)}</li>
                            <li><b>RAM:</b> {p.get('spec_ram', '—')}</li>
                            <li><b>Graphics:</b> {p.get('spec_gpu', '—')}</li>
                            <li><b>Weight:</b> {p.get('spec_weight', '—')}</li>
                        </ul>
                        <p style="font-size:0.82rem; opacity:0.85;"><b>Reason:</b> {p.get('match_reason', '')}</p>
                    </div>
                    """,
                    unsafe_allow_html=True
                )
        else:
            st.info("No ranked products found. Check Comparator Agent response.")

        st.markdown("</div>", unsafe_allow_html=True)

    # ---- Step 4: Sales Pitch Agent (ONLY HIGHLIGHTS + REASONS) ----
    with st.container():
        st.markdown(
            """
            <div class="step-card">
                <div class="step-header">
                    <div class="step-badge">4</div>
                    <div>
                        <div class="step-title">Sales Pitch Agent – HP Storytelling</div>
                        <div class="step-subtitle">Converts specs into customer-facing advantages & reasons to buy.</div>
                    </div>
                    <div class="step-status">✅ Completed</div>
                </div>
            """,
            unsafe_allow_html=True
        )

    # ===============================
    # 📄 PDF PREVIEW AND DOWNLOAD
    # ===============================
    if pdf_bytes:
        st.markdown("## 📄 PDF Preview & Download")

        pdf_base64 = base64.b64encode(pdf_bytes).decode()
        pdf_display = f'<iframe src="data:application/pdf;base64,{pdf_base64}" width="100%" height="500"></iframe>'
    
        st.markdown(pdf_display, unsafe_allow_html=True)

        st.markdown(
            f"""
            <div style="margin-top: 10px;">
                <a href="data:application/pdf;base64,{pdf_base64}" download="HP_Proposal_Preview.pdf">
                    <button class="hp-btn">
                        📄 Download PDF
                    </button>
                </a>
            </div>
            """,
            unsafe_allow_html=True
        )
    else:
        st.info("PDF format not generated or preview unavailable.")



# ---- Step 5: PPT Deck Generation ----
st.markdown("## 📁 Final Step – HP Deck Export")

col_left, col_right = st.columns([1.1, 1.9])

with col_left:
    # st.markdown(
    #     """
    #     <div class="step-card">
    #         <div class="step-header">
    #             <div class="step-badge">5</div>
    #             <div>
    #                 <div class="step-title">PPT Deck Generated</div>
    #                 <div class="step-subtitle">Your HP-branded proposal is ready to download & present.</div>
    #             </div>
    #             <div class="step-status">✅ Ready</div>
    #         </div>
    #     """,
    #     unsafe_allow_html=True
    # )

    st.markdown(
        """
        <div class="highlight-box">
            <h4>✅ Presentation Ready</h4>
            <p style="font-size:0.9rem;">
                The PowerPoint deck has been generated based on the ranked HP products and the sales pitch.
                You can now download it and share with the customer.
            </p>
        </div>
        """,
        unsafe_allow_html=True
    )

    dl_link = get_download_link(ppt_bytes)
    st.markdown(
        f"""
        <div class="download-btn">
            <a href="{dl_link}" download="HP_Proposal_Professional.pptx">
                <button class="hp-btn">
                    💾 Download HP Deck PPT
                </button>
            </a>
        </div>
        """,
        unsafe_allow_html=True
    )

    st.markdown("</div>", unsafe_allow_html=True)




# ===============================
# 🎯 Creativity Evaluation
# ===============================
if sales_pitch and ranked_products:
    st.markdown("## ✨ Creativity Evaluation Results")

    try:
        eval_url = f"{FASTAPI_BASE_URL}/evaluate-creativity-llm"
        payload = {
            "query": query,
            "sales_pitch": sales_pitch,
            "products": ranked_products[:5]
        }

        resp = requests.post(eval_url, json=payload, timeout=90)
        creativity_scores = resp.json()

        overall_score = creativity_scores.get("final_score", 0)

        with st.expander(f"🌟 Creativity Score — {overall_score}/10"):
            for key, item in creativity_scores.items():
                if key == "final_score":
                    continue
                score = item.get("score", 0)
                reason = item.get("reason", "")

                st.write(f"**{key.replace('_',' ').title()} — {score}/10**")
                st.progress(score / 10)
                st.caption(reason)

    except Exception as e:
        st.warning(f"Creativity evaluation failed: {e}")


